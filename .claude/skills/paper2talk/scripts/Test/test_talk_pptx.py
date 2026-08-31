"""
Offline tests for talk_pptx.py, the renderer that builds on the lab gabarit rather
than re-creating it: layout selection by tier and by name, the template's own sample
slides dropped from the copy, every block kind drawn, notes in the notes pane, the
audience font floor applied, and a block the renderer cannot draw raising instead of
being silently skipped.

python-pptx's packaged default template stands in for the gabarit, so the suite needs
no lab file and no PowerPoint.
"""
import base64
import json
import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import _fixtures as fx  # noqa: E402

import talk_model as tm  # noqa: E402
import talk_pptx as tp  # noqa: E402

# A 1x1 PNG, so the figure and equation-image branches have a real file to place.
PNG_1PX = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def blank_template(path):
    """A template with the default layouts, plus two sample slides to be dropped."""
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.save(path)
    return path


class TestLayoutSelection(unittest.TestCase):
    def setUp(self):
        from pptx import Presentation
        self.prs = Presentation()

    def test_each_tier_finds_a_layout(self):
        for tier in ("title", "content", "divider", "thanks", "backup"):
            self.assertIsNotNone(tp.pick_layout(self.prs, tier))

    def test_the_title_tier_prefers_a_title_layout(self):
        self.assertIn("title", tp.pick_layout(self.prs, "title").name.lower())

    def test_the_divider_tier_finds_the_section_layout(self):
        self.assertIn("section", tp.pick_layout(self.prs, "divider").name.lower())

    def test_an_explicit_name_wins(self):
        chosen = tp.pick_layout(self.prs, "content", "Blank")
        self.assertEqual(chosen.name, "Blank")

    def test_an_unknown_name_raises_and_lists_what_exists(self):
        with self.assertRaises(ValueError) as ctx:
            tp.pick_layout(self.prs, "content", "1- Contenu")
        self.assertIn("Blank", str(ctx.exception))

    def test_the_custom_gabarit_names_are_matched_first(self):
        # The 16:9 lab gabarit names its layouts "1- Contenu", "Page separatrice",
        # "Fin"; the patterns must reach them, not only the Office defaults.
        import re
        for name, tier in (("1- Contenu", "content"), ("1- Titre", "title"),
                           ("Page separatrice", "divider"), ("Fin", "thanks")):
            self.assertTrue(
                any(re.search(p, name, re.I) for p in tp.LAYOUT_PATTERNS[tier]),
                f"{name} unmatched for tier {tier}",
            )


class TestRender(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp()
        self.template = blank_template(os.path.join(self.dir, "gabarit.pptx"))
        self.fig = os.path.join(self.dir, "fig.png")
        with open(self.fig, "wb") as fh:
            fh.write(PNG_1PX)
        self.out = os.path.join(self.dir, "deck.pptx")

    def _model(self, blocks, **meta):
        return fx.model([
            {"n": 1, "kind": "title", "title": "T", "blocks": [], "notes": "opening"},
            {"n": 2, "kind": "content", "title": "A claim, not a topic",
             "blocks": blocks, "notes": "we discuss everything shown here"},
        ], **meta)

    def test_the_sample_slides_of_the_template_are_dropped(self):
        from pptx import Presentation

        self.assertEqual(len(Presentation(self.template).slides), 2)
        tp.render(self._model([{"kind": "figure", "asset": self.fig,
                                "keywords": ["fig"]}]),
                  self.template, self.out, verbose=False)
        prs = Presentation(self.out)
        self.assertEqual(len(prs.slides), 2)          # the model's two, not the samples

    def test_every_block_kind_draws(self):
        from pptx import Presentation

        blocks = [
            {"kind": "bullets", "items": ["one", "two"]},
            {"kind": "figure", "asset": self.fig, "caption": "Fig. 1", "keywords": ["fig"]},
            {"kind": "takeaway", "text": "it works", "support": "measured"},
            {"kind": "cards", "items": [{"title": "A", "text": "x"},
                                        {"title": "B", "text": "y"}]},
            {"kind": "chips", "items": [{"label": "low", "color": "2E7D32", "glyph": "L"}]},
            {"kind": "stats", "items": [{"value": "95.2 %", "label": "accuracy"}]},
            {"kind": "table", "rows": [["Model", "Accuracy"], ["SVM", "95.2 %"]],
             "keywords": ["accuracy"]},
            {"kind": "matrix", "rows": [[0.9, 0.1], [0.2, 0.8]], "keywords": ["matrix"]},
            {"kind": "zoneband", "zones": [{"fraction": 0.6, "label": "safe"},
                                           {"fraction": 0.4, "label": "alert"}],
             "keywords": ["zone"]},
            {"kind": "chart", "series": [{"name": "beta", "points": [[0, 0.1], [1, 0.7]]}],
             "min": 0, "max": 1, "keywords": ["curve"]},
            {"kind": "equation", "tex": "y = ax + b", "keywords": ["slope"]},
        ]
        self.assertEqual({b["kind"] for b in blocks}, set(tm.BLOCK_KINDS))
        tp.render(self._model(blocks), self.template, self.out, verbose=False)
        prs = Presentation(self.out)
        shapes = prs.slides[1].shapes
        self.assertTrue(any(s.shape_type is not None for s in shapes))
        self.assertTrue(any(s.has_table for s in shapes))
        self.assertTrue(any(s.has_chart for s in shapes))
        self.assertGreater(len(shapes), 10)

    def test_notes_land_in_the_notes_pane(self):
        from pptx import Presentation

        tp.render(self._model([{"kind": "figure", "asset": self.fig,
                                "keywords": ["fig"]}]),
                  self.template, self.out, verbose=False)
        prs = Presentation(self.out)
        text = prs.slides[1].notes_slide.notes_text_frame.text
        self.assertIn("we discuss everything shown here", text)

    def test_the_table_carries_the_house_style(self):
        from pptx import Presentation

        tp.render(self._model([{"kind": "table",
                                "rows": [["Model", "Accuracy"], ["SVM", "95.2 %"]],
                                "keywords": ["accuracy"]}]),
                  self.template, self.out, verbose=False)
        table = next(s.table for s in Presentation(self.out).slides[1].shapes if s.has_table)
        self.assertFalse(table.horz_banding)
        header = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        first_col = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
        body = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
        self.assertTrue(header.font.bold)
        self.assertTrue(first_col.font.bold)
        self.assertFalse(body.font.bold)

    def test_the_audience_font_floor_reaches_the_runs(self):
        from pptx import Presentation
        from pptx.util import Pt

        for audience, expected in (("field", 16), ("public", 20)):
            out = os.path.join(self.dir, f"deck_{audience}.pptx")
            tp.render(self._model([{"kind": "bullets", "items": ["one"]},
                                   {"kind": "figure", "asset": self.fig,
                                    "keywords": ["fig"]}],
                                  audience=audience),
                      self.template, out, verbose=False)
            sizes = []
            for shape in Presentation(out).slides[1].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        sizes += [r.font.size for r in para.runs if r.font.size]
            self.assertTrue(sizes)
            self.assertGreaterEqual(min(sizes), Pt(expected - 2))  # captions may go lower
            self.assertIn(Pt(expected), sizes)

    def test_a_block_the_renderer_cannot_draw_raises_before_anything_is_written(self):
        model = self._model([{"kind": "carousel", "items": []}])
        with self.assertRaises(tm.RendererGap):
            tp.render(model, self.template, self.out, verbose=False)
        self.assertFalse(os.path.exists(self.out))

    def test_a_missing_figure_is_named(self):
        model = self._model([{"kind": "figure", "asset": "nowhere.png",
                              "keywords": ["fig"]}])
        with self.assertRaises(FileNotFoundError):
            tp.render(model, self.template, self.out, verbose=False)

    def test_the_layout_map_overrides_the_tier_choice(self):
        from pptx import Presentation

        map_path = os.path.join(self.dir, "map.json")
        with open(map_path, "w", encoding="utf8") as fh:
            json.dump({"content": "Blank"}, fh)
        with open(map_path, encoding="utf8") as fh:
            layout_map = json.load(fh)
        tp.render(self._model([{"kind": "figure", "asset": self.fig,
                                "keywords": ["fig"]}]),
                  self.template, self.out, layout_map, verbose=False)
        prs = Presentation(self.out)
        self.assertEqual(prs.slides[1].slide_layout.name, "Blank")

    def test_the_gabarit_file_itself_is_never_written_to(self):
        before = os.path.getmtime(self.template)
        size = os.path.getsize(self.template)
        tp.render(self._model([{"kind": "figure", "asset": self.fig,
                                "keywords": ["fig"]}]),
                  self.template, self.out, verbose=False)
        self.assertEqual(os.path.getmtime(self.template), before)
        self.assertEqual(os.path.getsize(self.template), size)


class TestCli(unittest.TestCase):
    def setUp(self):
        self.dir = tempfile.mkdtemp()
        self.template = blank_template(os.path.join(self.dir, "gabarit.pptx"))

    def test_list_layouts_exits_zero(self):
        self.assertEqual(tp.main(["--list-layouts", "--template", self.template]), 0)

    def test_a_model_without_out_is_a_usage_error(self):
        model = os.path.join(self.dir, "m.json")
        with open(model, "w", encoding="utf8") as fh:
            json.dump(fx.model([]), fh)
        self.assertEqual(tp.main([model, "--template", self.template]), 2)


if __name__ == "__main__":
    unittest.main(verbosity=2)
