"""
talk_pptx - render a talk_model.json onto the lab PowerPoint gabarit.

The gabarit is a real template, not a background image: the 16:9 file carries 18
named layouts (`1- Titre`, `1- Contenu` ... `6- Contenu`, `Page separatrice`,
`Fin`), each with its own full-bleed brand picture and real placeholders, and the
4:3 file carries the standard eleven layouts over a branded master. A generator
that builds a deck from scratch cannot use any of that - pptxgenjs has no way to
read another file's layouts - so the PowerPoint target opens the gabarit itself and
adds slides from ITS layouts. Nothing about the branding is re-created.

    python talk_pptx.py <talk_model.json> --template <gabarit.pptx> --out <deck.pptx>
    python talk_pptx.py model.json --template g.pptx --out d.pptx --layout-map map.json
    python talk_pptx.py --list-layouts --template <gabarit.pptx>

The template's own sample slides are dropped from the copy; its masters, layouts,
theme and media are kept. The gabarit file itself is never written to.

`--layout-map` overrides the tier-to-layout choice, and a slide may name its layout
directly with a `layout` field:

    {"title": "1- Titre", "content": "1- Contenu", "divider": "Page separatrice",
     "thanks": "Fin", "backup": "Vide"}

Equations are the known weakness of this target: PowerPoint has no LaTeX, so a block
carrying `image` is placed as a picture and a bare `tex` is set as text. For an
equation-dense in-field talk, the Beamer target is the right answer.

Exit codes: 0 written, 1 a block the renderer cannot draw, 2 a usage or file error.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys

import talk_model as tm
import talk_rules as rules

TAG = "[PPTX]"

TFONT = "Cambria"   # titles, equations, stat numbers
BFONT = "Calibri"   # body

# How a slide tier finds its layout when the model and --layout-map say nothing.
# Ordered patterns, first match wins; the custom gabarit names are tried before the
# standard Office ones so a branded layout is preferred over a generic one.
LAYOUT_PATTERNS: dict[str, tuple[str, ...]] = {
    "title": (r"^1-\s*titre", r"titre.*diapositive", r"diapositive de titre",
              r"^titre", r"^title"),
    "content": (r"^1-\s*contenu", r"titre et contenu", r"contenu", r"title and content",
                r"^content"),
    "divider": (r"s[eé]paratrice", r"en-t[eê]te de section", r"section"),
    "thanks": (r"^fin$", r"merci", r"thank"),
    "backup": (r"^vide$", r"^blank$", r"^1-\s*contenu", r"contenu", r"content"),
}

# Semantic triad reused when a block names no colour of its own.
DEFAULT_PALETTE = {
    "brand": "006341", "ink": "23262A", "muted": "6B6F63",
    "card": "F1F2EC", "cardDk": "E4E7DA", "white": "FFFFFF",
    "s1": "2E7D32", "s2": "C57A00", "s3": "B3221F",
}


def _pptx():
    """Import python-pptx with an actionable message when it is missing."""
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.chart.data import CategoryChartData  # noqa: F401
        from pptx.enum.chart import XL_CHART_TYPE  # noqa: F401
        from pptx.enum.shapes import MSO_SHAPE  # noqa: F401
        from pptx.enum.text import PP_ALIGN  # noqa: F401
        from pptx.dml.color import RGBColor  # noqa: F401
        from pptx.util import Emu, Inches, Pt  # noqa: F401
    except ImportError as exc:  # pragma: no cover - environment-dependent
        raise RuntimeError(
            f"{TAG} python-pptx is required for the gabarit renderer: pip install python-pptx"
        ) from exc
    import pptx
    return pptx


def rgb(hex_or_name: str, palette: dict):
    """Resolve a palette key or a six-digit hex to an RGBColor."""
    from pptx.dml.color import RGBColor

    value = palette.get(hex_or_name, hex_or_name) if hex_or_name else palette["brand"]
    value = str(value).lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        value = palette["brand"]
    return RGBColor.from_string(value.upper())


def layouts_of(prs) -> list:
    """Every layout in the file, across masters, in a stable order."""
    seen = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            seen.append(layout)
    return seen


def pick_layout(prs, tier: str, wanted: str | None = None):
    """
    --------------------------------------------------------------------------
    Purpose:
        Choose the gabarit layout a slide is built on. An explicit name always
        wins; otherwise the tier's ordered patterns are matched against the
        layout names, so the branded `1- Contenu` is preferred over a generic
        `Titre et contenu` and both over nothing.

    Inputs:
        prs: the opened Presentation
        tier (str): content, title, thanks, divider, or backup
        wanted (str): an explicit layout name from the model or --layout-map

    Outputs:
        layout: the chosen SlideLayout
    --------------------------------------------------------------------------
    """
    available = layouts_of(prs)
    if wanted:
        for layout in available:
            if layout.name.strip().lower() == wanted.strip().lower():
                return layout
        raise ValueError(
            f"{TAG} layout {wanted!r} is not in the template; available: "
            + ", ".join(sorted(lay.name for lay in available))
        )
    for pattern in LAYOUT_PATTERNS.get(tier, ()):
        for layout in available:
            if re.search(pattern, layout.name.strip(), re.I):
                return layout
    return available[0]


def drop_sample_slides(prs) -> int:
    """
    --------------------------------------------------------------------------
    Purpose:
        Remove the slides the gabarit ships as examples, keeping its masters,
        layouts, theme and media. Six of them ride along in the 16:9 file, and
        shipping them inside a conference deck is not a subtle defect.

    Inputs:
        prs: the opened Presentation

    Outputs:
        n (int): how many slides were dropped
    --------------------------------------------------------------------------
    """
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    for entry in entries:
        prs.part.drop_rel(entry.rId)
        id_list.remove(entry)
    return len(entries)


def body_band(slide, prs, layout) -> tuple:
    """
    --------------------------------------------------------------------------
    Purpose:
        Find the rectangle a slide may draw in, in EMU. The body placeholder of
        the gabarit layout is the honest answer when there is one, because the
        lab already positioned it clear of the banner and the wordmark; without
        one, fall back to a margin inside the canvas.

    Inputs:
        slide: the new slide
        prs: the presentation (for the canvas size)
        layout: the layout the slide was built on

    Outputs:
        band (tuple): left, top, width, height in EMU
    --------------------------------------------------------------------------
    """
    from pptx.util import Inches

    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0 and shape.has_text_frame:
            return shape.left, shape.top, shape.width, shape.height
    for shape in layout.placeholders:
        if shape.placeholder_format.idx != 0 and shape.has_text_frame:
            return shape.left, shape.top, shape.width, shape.height
    margin = Inches(0.6)
    top = Inches(1.7)
    return (margin, top, prs.slide_width - 2 * margin,
            prs.slide_height - top - Inches(0.8))


def set_title(slide, text: str, font_pt: float, palette: dict) -> bool:
    """Fill the layout's title placeholder; returns False when it has none."""
    from pptx.util import Pt

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 or \
                str(shape.placeholder_format.type).startswith("TITLE"):
            shape.text_frame.text = text or ""
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_pt)
                    run.font.name = TFONT
                    run.font.bold = True
            return True
    return False


def _textbox(slide, left, top, width, height, text, size_pt, palette,
             bold=False, italic=False, color=None, align=None, font=BFONT):
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = str(text)
    for para in frame.paragraphs:
        if align is not None:
            para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.name = font
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(color or "ink", palette)
    return box


def render(model: dict, template: str, out_path: str,
           layout_map: dict | None = None, verbose: bool = True) -> str:
    """
    --------------------------------------------------------------------------
    Purpose:
        Build the deck on the gabarit: one slide per model entry, on the layout
        its tier selects, with the speaker notes in the notes pane.

    Inputs:
        model (dict): a parsed talk_model.json
        template (str): the gabarit .pptx (opened, never written to)
        out_path (str): destination deck
        layout_map (dict): tier -> layout name overrides
        verbose (bool): print the layout chosen for each slide

    Outputs:
        out_path (str): the file written
    --------------------------------------------------------------------------
    """
    _pptx()
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    # The renderer implements every block kind the model knows; the check is here
    # so a gap raises before a single slide exists rather than leaving a hole.
    tm.assert_renderable(model, set(tm.BLOCK_KINDS), "talk_pptx")

    meta = model.get("meta", {})
    audience = meta.get("audience", "field")
    profile = rules.audience_profile(audience)
    floor = profile["font_floor_pt"]
    cap = profile["caption_floor_pt"]
    title_pt = floor + (6 if audience == "public" else 8)
    palette = dict(DEFAULT_PALETTE)
    palette.update(model.get("palette") or {})
    layout_map = layout_map or {}

    prs = Presentation(template)
    dropped = drop_sample_slides(prs)
    if verbose and dropped:
        print(f"{TAG} dropped {dropped} sample slide(s) from the gabarit copy")

    canvas_in = (prs.slide_width / 914400, prs.slide_height / 914400)
    if verbose:
        print(f"{TAG} canvas {canvas_in[0]:.3f} x {canvas_in[1]:.3f} in "
              f"({len(layouts_of(prs))} layouts available)")

    for entry in model.get("slides", []):
        tier = tm.slide_tier(entry)
        wanted = entry.get("layout") or layout_map.get(tier)
        layout = pick_layout(prs, tier, wanted)
        slide = prs.slides.add_slide(layout)
        set_title(slide, entry.get("title", ""), title_pt, palette)

        left, top, width, height = body_band(slide, prs, layout)
        cursor = top
        bottom = top + height

        def remaining():
            return max(Inches(0.4), bottom - cursor)

        for block in entry.get("blocks", []):
            kind = block.get("kind")
            if kind == "bullets":
                items = block["items"]
                box = slide.shapes.add_textbox(left, cursor, width,
                                               min(remaining(), Inches(0.42 * len(items) + 0.2)))
                frame = box.text_frame
                frame.word_wrap = True
                for i, item in enumerate(items):
                    para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                    para.text = str(item)
                    para.level = int(block.get("level", 0))
                    for run in para.runs:
                        run.font.size = Pt(floor)
                        run.font.name = BFONT
                        run.font.color.rgb = rgb("ink", palette)
                cursor += Inches(0.42 * len(items) + 0.3)

            elif kind == "figure":
                asset = block["asset"]
                if not os.path.exists(asset):
                    raise FileNotFoundError(f"{TAG} figure {asset} not found")
                pic_h = Inches(float(block.get("h", 2.8)))
                pic = slide.shapes.add_picture(asset, left, cursor, height=pic_h)
                if pic.width > width:      # keep it inside the band, uniformly scaled
                    scale = width / pic.width
                    pic.width = int(pic.width * scale)
                    pic.height = int(pic.height * scale)
                pic.left = left + int((width - pic.width) / 2)
                cursor += pic.height + Inches(0.06)
                label = " ".join(x for x in (block.get("caption"), block.get("cite")) if x)
                if label:
                    _textbox(slide, left, cursor, width, Inches(0.3), label, cap,
                             palette, italic=True, color="muted")
                    cursor += Inches(0.34)

            elif kind == "takeaway":
                band_h = Inches(0.95 if block.get("support") else 0.62)
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               left, cursor, width, band_h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb("cardDk", palette)
                shape.line.color.rgb = rgb("cardDk", palette)
                shape.text_frame.word_wrap = True
                shape.text_frame.text = str(block["text"])
                for run in shape.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(floor + 1)
                    run.font.name = TFONT
                    run.font.bold = True
                    run.font.color.rgb = rgb("brand", palette)
                if block.get("support"):
                    para = shape.text_frame.add_paragraph()
                    para.text = str(block["support"])
                    for run in para.runs:
                        run.font.size = Pt(floor)
                        run.font.name = BFONT
                        run.font.color.rgb = rgb("ink", palette)
                cursor += band_h + Inches(0.12)

            elif kind == "cards":
                items = block["items"]
                gap = Inches(0.16)
                card_w = int((width - gap * (len(items) - 1)) / len(items))
                card_h = Inches(float(block.get("h", 1.5)))
                for i, item in enumerate(items):
                    x = left + i * (card_w + gap)
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   x, cursor, card_w, card_h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = rgb("card", palette)
                    shape.line.color.rgb = rgb("card", palette)
                    frame = shape.text_frame
                    frame.word_wrap = True
                    frame.text = str(item.get("title", ""))
                    for run in frame.paragraphs[0].runs:
                        run.font.size = Pt(floor)
                        run.font.bold = True
                        run.font.name = BFONT
                        run.font.color.rgb = rgb("brand", palette)
                    para = frame.add_paragraph()
                    para.text = str(item.get("text", ""))
                    for run in para.runs:
                        run.font.size = Pt(floor)
                        run.font.name = BFONT
                        run.font.color.rgb = rgb("ink", palette)
                cursor += card_h + Inches(0.14)

            elif kind == "chips":
                # Semantic chips: the colour IS the class. A chip whose colour
                # means nothing is ornament and should be dropped instead.
                items = block["items"]
                dia = Inches(0.4)
                step = int(width / max(1, len(items)))
                for i, item in enumerate(items):
                    x = left + i * step
                    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cursor, dia, dia)
                    disc.fill.solid()
                    disc.fill.fore_color.rgb = rgb(item.get("color", "brand"), palette)
                    disc.line.color.rgb = rgb(item.get("color", "brand"), palette)
                    disc.text_frame.text = str(item.get("glyph", ""))
                    for run in disc.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(floor)
                        run.font.bold = True
                        run.font.name = TFONT
                        run.font.color.rgb = rgb("white", palette)
                    _textbox(slide, x + dia + Inches(0.1), cursor,
                             step - dia - Inches(0.15), dia,
                             item.get("label", ""), floor, palette, bold=True)
                cursor += dia + Inches(0.2)

            elif kind == "stats":
                items = block["items"]
                col = int(width / len(items))
                for i, item in enumerate(items):
                    x = left + i * col
                    _textbox(slide, x, cursor, col, Inches(0.6), item.get("value", ""),
                             floor + 12, palette, bold=True, color="brand",
                             align=PP_ALIGN.CENTER, font=TFONT)
                    _textbox(slide, x, cursor + Inches(0.6), col, Inches(0.4),
                             item.get("label", ""), floor, palette,
                             align=PP_ALIGN.CENTER)
                cursor += Inches(1.05)

            elif kind == "table":
                rows = block["rows"]
                n_rows, n_cols = len(rows), len(rows[0])
                # A row grows with its content, so this height is a floor and the
                # rendered page still has to be looked at.
                table_h = Inches(0.38 * n_rows)
                frame = slide.shapes.add_table(n_rows, n_cols, left, cursor,
                                               width, table_h).table
                # House style (.claude/CLAUDE.md): bold header row on a 10 % grey,
                # bold first column, no banding. The default table style paints an
                # accent-blue header that belongs to no lab palette.
                frame.first_row = False
                frame.horz_banding = False
                for r, row in enumerate(rows):
                    for c, value in enumerate(row):
                        cell = frame.cell(r, c)
                        cell.text = str(value)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = rgb(
                            "cardDk" if r == 0 else "white", palette)
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(floor)
                                run.font.name = BFONT
                                run.font.bold = (r == 0 or c == 0)
                                run.font.color.rgb = rgb("ink", palette)
                cursor += table_h + Inches(0.14)

            elif kind == "matrix":
                rows = block["rows"]
                cell_side = Inches(float(block.get("cell", 0.62)))
                for r, row in enumerate(rows):
                    for c, value in enumerate(row):
                        t = max(0.0, min(1.0, float(value)))
                        box = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, left + c * cell_side,
                            cursor + r * cell_side, cell_side, cell_side)
                        box.fill.solid()
                        box.fill.fore_color.rgb = rgb("brand", palette)
                        box.fill.transparency = 1.0 - t
                        box.line.color.rgb = rgb("white", palette)
                        box.text_frame.text = str(value)
                        for run in box.text_frame.paragraphs[0].runs:
                            run.font.size = Pt(floor)
                            run.font.name = BFONT
                            run.font.color.rgb = rgb(
                                "white" if t > 0.55 else "ink", palette)
                cursor += cell_side * len(rows) + Inches(0.2)

            elif kind == "zoneband":
                band_h = Inches(float(block.get("h", 0.5)))
                x = left
                for zone in block["zones"]:
                    zw = int(width * float(zone.get("fraction", 0)))
                    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cursor,
                                                 zw, band_h)
                    box.fill.solid()
                    box.fill.fore_color.rgb = rgb(zone.get("color", "brand"), palette)
                    box.line.color.rgb = rgb("white", palette)
                    box.text_frame.text = str(zone.get("label", ""))
                    for run in box.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(floor)
                        run.font.name = BFONT
                        run.font.bold = True
                        run.font.color.rgb = rgb("white", palette)
                    x += zw
                cursor += band_h + Inches(0.24)

            elif kind == "chart":
                data = CategoryChartData()
                series = block["series"]
                first = series[0]
                data.categories = [str(p[0]) for p in first["points"]]
                for serie in series:
                    data.add_series(serie.get("name", ""),
                                    [float(p[1]) for p in serie["points"]])
                chart_h = Inches(float(block.get("h", 2.8)))
                frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, cursor,
                                               width, chart_h, data)
                chart = frame.chart
                # Declared limits, no marker, no category gridline - and never a
                # secondary value axis, which PowerPoint rejects.
                chart.value_axis.minimum_scale = float(block.get("min", 0))
                chart.value_axis.maximum_scale = float(block.get("max", 1))
                chart.category_axis.has_major_gridlines = False
                chart.has_legend = block.get("legend", True)
                cursor += chart_h + Inches(0.16)

            elif kind == "equation":
                # No LaTeX in PowerPoint. An exported image is faithful; the text
                # fallback is not, which is why an equation-dense talk belongs on
                # the Beamer target.
                if block.get("image") and os.path.exists(block["image"]):
                    pic = slide.shapes.add_picture(
                        block["image"], left, cursor,
                        height=Inches(float(block.get("h", 0.6))))
                    cursor += pic.height + Inches(0.12)
                else:
                    _textbox(slide, left, cursor, width, Inches(0.5),
                             block.get("plain") or block["tex"], floor + 2, palette,
                             align=PP_ALIGN.CENTER, font=TFONT)
                    cursor += Inches(0.6)
                if block.get("where"):
                    _textbox(slide, left, cursor, width, Inches(0.3), block["where"],
                             cap, palette, italic=True, color="muted")
                    cursor += Inches(0.34)

        if cursor > bottom and verbose:
            print(f"{TAG} slide {entry.get('n')} runs "
                  f"{Emu(int(cursor - bottom)).inches:.2f} in past its content band; "
                  "inspect the rendered page")

        # Notes live in the notes pane, never in a text box on the slide.
        if entry.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(entry["notes"])

    os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
    prs.save(out_path)
    if verbose:
        print(f"{TAG} wrote {out_path} ({len(model.get('slides', []))} slides)")
    return out_path


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Render a talk model onto a PowerPoint gabarit.")
    ap.add_argument("model", nargs="?")
    ap.add_argument("--template", required=True, help="the lab gabarit .pptx")
    ap.add_argument("--out", help="destination deck")
    ap.add_argument("--layout-map", help="JSON mapping tier -> layout name")
    ap.add_argument("--list-layouts", action="store_true",
                    help="print the template's layout names and exit")
    args = ap.parse_args(argv)

    try:
        _pptx()
        from pptx import Presentation
    except RuntimeError as exc:
        print(str(exc), file=sys.stderr)
        return 2

    if args.list_layouts:
        try:
            prs = Presentation(args.template)
        except Exception as exc:
            print(f"{TAG} cannot open {args.template}: {exc}", file=sys.stderr)
            return 2
        print(f"{TAG} {prs.slide_width / 914400:.3f} x {prs.slide_height / 914400:.3f} in")
        for layout in layouts_of(prs):
            print(f"{TAG}   {layout.name}")
        return 0

    if not args.model or not args.out:
        print(f"{TAG} give a model and --out (or use --list-layouts)", file=sys.stderr)
        return 2

    try:
        model = tm.load(args.model)
        layout_map = None
        if args.layout_map:
            with open(args.layout_map, encoding="utf8") as fh:
                layout_map = json.load(fh)
        render(model, args.template, args.out, layout_map)
    except tm.RendererGap as exc:
        print(str(exc), file=sys.stderr)
        return 1
    except (OSError, ValueError, RuntimeError) as exc:
        print(f"{TAG} {exc}", file=sys.stderr)
        return 2
    return 0


if __name__ == "__main__":
    sys.exit(main())
